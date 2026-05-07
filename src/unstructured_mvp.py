from __future__ import annotations

import re
from collections import Counter
from pathlib import Path

import pandas as pd

KEYWORD_THEMES = {
    "engagement": ["engagement", "satisfaction", "enquete", "feedback", "qvt", "wellbeing"],
    "diversity_inclusion": ["diversite", "inclusion", "mixite", "equite", "egalite"],
    "training": ["formation", "training", "learning", "competence", "certification"],
    "retention_mobility": ["turnover", "retention", "mobilite", "attrition", "absenteisme"],
    "governance": ["governance", "ethique", "compliance", "manager", "leadership"],
}

STOPWORDS = {
    "de",
    "la",
    "le",
    "les",
    "des",
    "et",
    "en",
    "du",
    "au",
    "aux",
    "a",
    "d",
    "l",
    "pour",
    "sur",
    "avec",
    "the",
    "and",
    "to",
    "of",
    "in",
    "for",
}


def _safe_read_pdf(path: Path) -> str:
    try:
        import pypdf  # type: ignore
    except Exception:
        return ""
    text_chunks: list[str] = []
    try:
        reader = pypdf.PdfReader(str(path))
        for page in reader.pages:
            text_chunks.append(page.extract_text() or "")
    except Exception:
        return ""
    return "\n".join(text_chunks)


def _safe_read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return ""
    chunks: list[str] = []
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    chunks.append(text)
    except Exception:
        return ""
    return "\n".join(chunks)


def discover_documents(data_root: Path) -> list[Path]:
    return sorted([*data_root.rglob("*.pdf"), *data_root.rglob("*.pptx")])


def extract_unstructured_corpus(data_root: Path) -> pd.DataFrame:
    rows = []
    for file_path in discover_documents(data_root):
        if file_path.suffix.lower() == ".pdf":
            text = _safe_read_pdf(file_path)
            source_type = "pdf"
        else:
            text = _safe_read_pptx(file_path)
            source_type = "pptx"
        rows.append(
            {
                "file_name": file_path.name,
                "relative_path": str(file_path.relative_to(data_root)),
                "source_type": source_type,
                "char_count": len(text),
                "word_count": len(re.findall(r"[A-Za-zÀ-ÿ]{3,}", text)),
                "text": text,
            }
        )
    return pd.DataFrame(rows)


def _tokenize(text: str) -> list[str]:
    tokens = [t.lower() for t in re.findall(r"[A-Za-zÀ-ÿ]{4,}", text or "")]
    return [t for t in tokens if t not in STOPWORDS]


def compute_theme_signals(corpus: pd.DataFrame) -> pd.DataFrame:
    if corpus.empty:
        return pd.DataFrame(columns=["theme", "hit_count", "normalized_score"])
    all_text = " ".join(corpus["text"].fillna("").tolist()).lower()
    totals = []
    total_words = max(sum(corpus["word_count"].fillna(0)), 1)
    for theme, keywords in KEYWORD_THEMES.items():
        pattern = r"|".join(re.escape(k) for k in keywords)
        hits = len(re.findall(pattern, all_text))
        totals.append({"theme": theme, "hit_count": hits, "normalized_score": round((hits / total_words) * 1000, 4)})
    return pd.DataFrame(totals).sort_values("normalized_score", ascending=False)


def compute_top_terms(corpus: pd.DataFrame, top_n: int = 20) -> pd.DataFrame:
    if corpus.empty:
        return pd.DataFrame(columns=["term", "count"])
    token_counter = Counter()
    for txt in corpus["text"].fillna(""):
        token_counter.update(_tokenize(txt))
    rows = [{"term": term, "count": count} for term, count in token_counter.most_common(top_n)]
    return pd.DataFrame(rows)


def build_unstructured_summary(corpus: pd.DataFrame) -> dict[str, float | int]:
    if corpus.empty:
        return {"documents": 0, "pdf_files": 0, "pptx_files": 0, "extractable_docs_pct": 0.0, "total_words": 0}
    extractable = corpus["word_count"].fillna(0).gt(0)
    return {
        "documents": int(len(corpus)),
        "pdf_files": int(corpus["source_type"].eq("pdf").sum()),
        "pptx_files": int(corpus["source_type"].eq("pptx").sum()),
        "extractable_docs_pct": round(100 * float(extractable.mean()), 1),
        "total_words": int(corpus["word_count"].sum()),
    }
